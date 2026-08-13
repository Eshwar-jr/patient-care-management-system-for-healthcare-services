import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# Color Palette
BG_COLOR = RGBColor(248, 250, 252)        # Light Slate / Gray background (#F8FAFC)
CARD_BG = RGBColor(255, 255, 255)         # Pure White card (#FFFFFF)
CARD_BORDER = RGBColor(203, 213, 225)     # Subtle Gray border (#CBD5E1)
PRIMARY_NAVY = RGBColor(15, 23, 42)       # Dark Navy (#0F172A)
PRIMARY_BLUE = RGBColor(2, 132, 199)      # Vibrant Ocean Blue (#0284C7)
TEAL_ACCENT = RGBColor(13, 148, 136)      # Medical Teal (#0D9488)
TEXT_DARK = RGBColor(15, 23, 42)          # Primary Text (#0F172A)
TEXT_MUTED = RGBColor(71, 85, 105)        # Muted Slate Text (#475569)
TEXT_LIGHT = RGBColor(255, 255, 255)      # White Text
HEADER_BG = RGBColor(15, 23, 42)          # Header Dark Slate
BADGE_BG = RGBColor(224, 242, 254)        # Light Blue Badge (#E0F2FE)
BADGE_TEXT = RGBColor(3, 105, 161)        # Dark Blue Badge Text (#0369A1)
PLACEHOLDER_BG = RGBColor(241, 245, 249)  # Light Slate (#F1F5F9)
PLACEHOLDER_BORDER = RGBColor(148, 163, 184) # Slate Border (#94A3B8)
SUCCESS_GREEN = RGBColor(22, 163, 74)     # Success Green (#16A34A)
CARD_HDR_BG = RGBColor(241, 245, 249)     # Card Header Light Slate (#F1F5F9)

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"

def set_slide_background(slide, color=BG_COLOR):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_header(slide, badge_text, title_text, subtitle_text=""):
    # Category / Badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.8), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = BADGE_BG
    badge.line.color.rgb = RGBColor(186, 230, 253)
    
    tf_b = badge.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = badge_text.upper()
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.name = FONT_HEADING
    p_b.font.color.rgb = BADGE_TEXT
    p_b.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.733), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = FONT_HEADING
    p.font.color.rgb = PRIMARY_NAVY
    
    if subtitle_text:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(12)
        p_sub.font.name = FONT_BODY
        p_sub.font.color.rgb = TEXT_MUTED

def add_footer(slide, current_slide, total_slides=13):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = CARD_BORDER
    line.line.fill.background()
    
    foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(8.0), Inches(0.35))
    tf = foot_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Integrated Patient Care Management System (IPCMS) | College Project Evaluation"
    p.font.size = Pt(10)
    p.font.name = FONT_BODY
    p.font.color.rgb = TEXT_MUTED
    
    num_box = slide.shapes.add_textbox(Inches(10.533), Inches(7.05), Inches(2.0), Inches(0.35))
    tf_n = num_box.text_frame
    p_n = tf_n.paragraphs[0]
    p_n.text = f"Slide {current_slide} of {total_slides}"
    p_n.font.size = Pt(10)
    p_n.font.bold = True
    p_n.font.name = FONT_BODY
    p_n.font.color.rgb = TEXT_MUTED
    p_n.alignment = PP_ALIGN.RIGHT

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card

def add_placeholder_frame(slide, left, top, width, height, label_text):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = PLACEHOLDER_BG
    card.line.color.rgb = PLACEHOLDER_BORDER
    card.line.width = Pt(1.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.1)
    
    p1 = tf.paragraphs[0]
    p1.text = "INSERT SCREENSHOT HERE"
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.name = FONT_HEADING
    p1.font.color.rgb = PRIMARY_BLUE
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = f"[{label_text}]"
    p2.font.size = Pt(10)
    p2.font.bold = True
    p2.font.name = FONT_BODY
    p2.font.color.rgb = TEXT_DARK
    p2.alignment = PP_ALIGN.CENTER

# ==============================================================================
# SLIDE 1: TITLE SLIDE
# ==============================================================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1, PRIMARY_NAVY)

# Title Card Accent Box
title_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1))
title_card.fill.solid()
title_card.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
title_card.line.color.rgb = RGBColor(51, 65, 85)
title_card.line.width = Pt(1.5)

tf1 = title_card.text_frame
tf1.word_wrap = True
tf1.margin_left = Inches(0.6)
tf1.margin_top = Inches(0.5)

# Badge
p = tf1.paragraphs[0]
p.text = "PROJECT EVALUATION PRESENTATION"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE
p.font.name = FONT_HEADING
p.alignment = PP_ALIGN.LEFT

# Main Title
p = tf1.add_paragraph()
p.text = "Integrated Patient Care Management System"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.font.name = FONT_HEADING

# Subtitle
p = tf1.add_paragraph()
p.text = "IPCMS - Hospital Management and Patient Care Platform"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
p.font.name = FONT_BODY

# Divider line
p = tf1.add_paragraph()
p.text = "---------------------------------------------"
p.font.size = Pt(14)
p.font.color.rgb = PRIMARY_BLUE

# Presenter & Guide Placeholders Box
place_box = slide1.shapes.add_textbox(Inches(1.6), Inches(4.3), Inches(10.133), Inches(1.6))
tf_p = place_box.text_frame
tf_p.word_wrap = True

p = tf_p.paragraphs[0]
p.text = "PRESENTER PLACEHOLDER:"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE
p.font.name = FONT_HEADING

p = tf_p.add_paragraph()
p.text = "• Student Name: [Insert Presenter Name]  |  Register No: [Insert Registration Number]"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_LIGHT
p.font.name = FONT_BODY

p = tf_p.add_paragraph()
p.text = "GUIDE / MENTOR PLACEHOLDER:"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE
p.font.name = FONT_HEADING

p = tf_p.add_paragraph()
p.text = "• Guide Name: [Insert Guide / Mentor Name & Designation]  |  Department: Computer Science & Engineering"
p.font.size = Pt(13)
p.font.color.rgb = TEXT_LIGHT
p.font.name = FONT_BODY


# ==============================================================================
# SLIDE 2: PROBLEM STATEMENT & OBJECTIVES
# ==============================================================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_header(slide2, "PROJECT OVERVIEW", "Problem Statement & Project Objectives", "Identifying healthcare workflow challenges and defining system targets")
add_footer(slide2, 2)

# Left Box: Problem Statement
add_card(slide2, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.2))
box1 = slide2.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.3), Inches(5.0))
tf2_1 = box1.text_frame
tf2_1.word_wrap = True

p = tf2_1.paragraphs[0]
p.text = "PROBLEM STATEMENT"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

p = tf2_1.add_paragraph()
p.text = "Traditional manual processes in hospital environments lead to:"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED

problems = [
    "Fragmented Patient Information: Disconnected physical files and isolated records create data silos.",
    "Inefficient Appointment Scheduling: Manual booking causes scheduling conflicts and extended waiting times.",
    "Disconnected Consultations: Clinical findings are not linked automatically with patient medical history.",
    "Manual Prescriptions & Lab Workflows: Hand-written notes delay pharmacy fulfillment and lab result delivery.",
    "Opaque Billing Processes: Manual ledger entries increase error rates and hinder payment tracking.",
    "Lack of Centralized Analytics: Absence of real-time insights into disease trends and hospital revenue.",
    "Unmonitored Feedback: Inability to systematically capture and analyze patient satisfaction ratings."
]
for item in problems:
    p = tf2_1.add_paragraph()
    title, desc = item.split(": ", 1)
    run1 = p.add_run()
    run1.text = "• " + title + ": "
    run1.font.bold = True
    run1.font.size = Pt(11)
    run1.font.color.rgb = TEXT_DARK
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(11)
    run2.font.color.rgb = TEXT_MUTED

# Right Box: Objectives
add_card(slide2, Inches(6.8), Inches(1.5), Inches(5.733), Inches(5.2))
box2 = slide2.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.333), Inches(5.0))
tf2_2 = box2.text_frame
tf2_2.word_wrap = True

p = tf2_2.paragraphs[0]
p.text = "SYSTEM OBJECTIVES"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

p = tf2_2.add_paragraph()
p.text = "Core engineering goals of the IPCMS application:"
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED

objectives = [
    "Centralize Patient Data: Store comprehensive demographic and medical data securely in MySQL.",
    "Streamline Appointments: Enable online scheduling, doctor selection, and status updates.",
    "Electronic Health Records (EHR): Maintain full vital signs history, allergies, and diagnoses.",
    "Automate Prescriptions & Labs: Enable digital prescription writing and lab test tracking.",
    "Integrated Billing & Alerts: Auto-generate bills from care services and issue real-time alerts.",
    "Administrative Reporting: Provide customizable reports with PDF, Excel, and CSV export capabilities.",
    "RESTful API Architecture: Expose authenticated endpoints for external integration and data access.",
    "Testing & Security Hardening: Achieve 100% test pass rate with RBAC and CSRF/SQLi protection."
]
for item in objectives:
    p = tf2_2.add_paragraph()
    title, desc = item.split(": ", 1)
    run1 = p.add_run()
    run1.text = "• " + title + ": "
    run1.font.bold = True
    run1.font.size = Pt(11)
    run1.font.color.rgb = TEXT_DARK
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(11)
    run2.font.color.rgb = TEXT_MUTED


# ==============================================================================
# SLIDE 3: TECHNOLOGY STACK + SYSTEM ARCHITECTURE
# ==============================================================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_header(slide3, "ARCHITECTURE & TECH STACK", "Technology Stack & Layered System Architecture", "Verified technology stack and multi-layered web application design")
add_footer(slide3, 3)

# Left Column: Technology Stack
add_card(slide3, Inches(0.8), Inches(1.5), Inches(5.0), Inches(5.2))
box3_1 = slide3.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(4.6), Inches(5.0))
tf3_1 = box3_1.text_frame
tf3_1.word_wrap = True

p = tf3_1.paragraphs[0]
p.text = "VERIFIED TECHNOLOGY STACK"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

tech_list = [
    ("Core Language", "Python 3.10+"),
    ("Web Framework", "Flask 3.1.3"),
    ("ORM & Database", "Flask-SQLAlchemy 3.1, MySQL 8.0, PyMySQL"),
    ("Authentication", "Flask-Login 0.6, Werkzeug 3.1 (Bcrypt Hashing)"),
    ("Form Validation", "WTForms 3.2, Flask-WTF 1.3"),
    ("Frontend UI", "HTML5, CSS3, JavaScript, Bootstrap 5"),
    ("Data Visualization", "Chart.js (Interactive Dashboards)"),
    ("Document Export", "ReportLab (PDF), XlsxWriter (Excel), CSV"),
    ("WSGI Server & Env", "Waitress 3.0, python-dotenv"),
    ("Testing Suite", "Python unittest (Automated Suite)")
]

for category, tech in tech_list:
    p = tf3_1.add_paragraph()
    run1 = p.add_run()
    run1.text = f"• {category}: "
    run1.font.bold = True
    run1.font.size = Pt(11)
    run1.font.color.rgb = TEXT_DARK
    run2 = p.add_run()
    run2.text = tech
    run2.font.size = Pt(11)
    run2.font.color.rgb = PRIMARY_BLUE

# Right Column: System Architecture Diagram
add_card(slide3, Inches(6.1), Inches(1.5), Inches(6.433), Inches(5.2))
box3_2 = slide3.shapes.add_textbox(Inches(6.3), Inches(1.6), Inches(6.033), Inches(5.0))
tf3_2 = box3_2.text_frame
tf3_2.word_wrap = True

p = tf3_2.paragraphs[0]
p.text = "LAYERED SYSTEM ARCHITECTURE"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

layers = [
    ("USER ROLES", "Admin | Doctor | Nurse | Patient | Pharmacist", PRIMARY_BLUE),
    ("FRONTEND / TEMPLATES", "Jinja2 Templates + Bootstrap 5 + Chart.js", TEAL_ACCENT),
    ("ROUTES / CONTROLLERS", "Flask Blueprints / Routes (auth.py, dashboard.py, api.py)", PRIMARY_NAVY),
    ("BUSINESS LOGIC & SECURITY", "WTForms Validation + Flask-Login + RBAC Middleware", PRIMARY_BLUE),
    ("DATA ORM LAYER", "SQLAlchemy ORM (Models & Relationships)", TEAL_ACCENT),
    ("DATABASE PERSISTENCE", "MySQL Relational Database Engine", PRIMARY_NAVY)
]

y_pos = Inches(2.1)
for layer_title, layer_desc, color in layers:
    card = add_card(slide3, Inches(6.4), y_pos, Inches(5.8), Inches(0.55), bg_color=CARD_BG, border_color=color)
    tf_c = card.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = Inches(0.15)
    tf_c.margin_top = Inches(0.05)
    
    p = tf_c.paragraphs[0]
    p.text = layer_title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    
    p_sub = tf_c.add_paragraph()
    p_sub.text = layer_desc
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_DARK
    p_sub.alignment = PP_ALIGN.LEFT
    
    y_pos += Inches(0.68)

# Add supporting modules box at bottom right
supp_card = add_card(slide3, Inches(6.4), Inches(6.0), Inches(5.8), Inches(0.55), bg_color=BADGE_BG, border_color=BADGE_TEXT)
tf_s = supp_card.text_frame
tf_s.margin_left = Inches(0.15)
tf_s.margin_top = Inches(0.05)
p_s = tf_s.paragraphs[0]
p_s.text = "INTEGRATED SYSTEM SUBSYSTEMS"
p_s.font.size = Pt(9)
p_s.font.bold = True
p_s.font.color.rgb = BADGE_TEXT

p_s2 = tf_s.add_paragraph()
p_s2.text = "REST API Engine  |  Report Generator (PDF/Excel/CSV)  |  Analytics & Caching  |  Notification System"
p_s2.font.size = Pt(9)
p_s2.font.color.rgb = TEXT_DARK


# ==============================================================================
# SLIDE 4: MILESTONE 1 — DATABASE & AUTHENTICATION
# ==============================================================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_header(slide4, "MILESTONE 1 (1/2)", "Milestone 1 - Database Design & User Authentication", "Foundation setup, relational database configuration, user models, and role-based login portals")
add_footer(slide4, 4)

# Left Column: Features
add_card(slide4, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2))
box4 = slide4.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.1), Inches(5.0))
tf4 = box4.text_frame
tf4.word_wrap = True

p = tf4.paragraphs[0]
p.text = "IMPLEMENTED M1 WORK - DATABASE & AUTH"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

m1_1_items = [
    "Flask App Architecture: Core initialization with modular blueprints and extension bindings (`db`, `login_manager`).",
    "MySQL Database Connection: Configured PyMySQL driver with SQLAlchemy URI string support.",
    "User Persistence Model: Designed `User` model with attributes for full name, username, email, phone, role, and timestamp.",
    "Password Hashing: Integrated Werkzeug security generating salted password hashes for secure storage.",
    "Role-Based Authentication: Built 5 distinct user portals: Admin, Doctor, Nurse, Patient, and Pharmacist.",
    "Wrong-Portal Detection: Automatic detection preventing users from logging into unauthorized portal roles.",
    "Audit Activity Logging: Persistence of login/logout activities via `LoginActivity` table with IP tracking.",
    "Auto Schema Initialization: Automatic table creation on startup (`db.create_all()`)."
]

for item in m1_1_items:
    p = tf4.add_paragraph()
    title, desc = item.split(": ", 1)
    run1 = p.add_run()
    run1.text = "• " + title + ": "
    run1.font.bold = True
    run1.font.size = Pt(10.5)
    run1.font.color.rgb = TEXT_DARK
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(10.5)
    run2.font.color.rgb = TEXT_MUTED

# Right Column: 2 Screenshot Placeholders
add_placeholder_frame(slide4, Inches(6.6), Inches(1.5), Inches(5.933), Inches(3.1), "INSERT SCREENSHOT - Registration/Login Page")
add_placeholder_frame(slide4, Inches(6.6), Inches(4.8), Inches(5.933), Inches(1.9), "INSERT SCREENSHOT - Database/MySQL Tables")


# ==============================================================================
# SLIDE 5: MILESTONE 1 — CORE PATIENT MANAGEMENT
# ==============================================================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_header(slide5, "MILESTONE 1 (2/2)", "Milestone 1 - Patient & Core Hospital Management", "Patient onboarding, record persistence, appointment creation, and base dashboard navigation")
add_footer(slide5, 5)

# Left Column: Features & Workflow
add_card(slide5, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2))
box5 = slide5.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.1), Inches(5.0))
tf5 = box5.text_frame
tf5.word_wrap = True

p = tf5.paragraphs[0]
p.text = "PATIENT MANAGEMENT & CORE FLOW"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

m1_2_items = [
    "Patient Registration: Capturing full name, age, gender, phone, address, blood group, disease, Aadhaar, and email.",
    "Patient Records Listing: Filterable tabular view displaying registered patients with instant search support.",
    "WTForms Validation: Server-side field validation ensuring Aadhaar uniqueness and proper phone format.",
    "Appointment Creation: Core booking feature connecting registered patients with doctors, date, time, and reason.",
    "Role-Based Dashboards: Custom landing pages per role with personalized navigation and quick action controls.",
    "MySQL Entity Persistence: Immediate relational saving across `User`, `Patient`, and `Appointment` models."
]

for item in m1_2_items:
    p = tf5.add_paragraph()
    title, desc = item.split(": ", 1)
    run1 = p.add_run()
    run1.text = "• " + title + ": "
    run1.font.bold = True
    run1.font.size = Pt(10.5)
    run1.font.color.rgb = TEXT_DARK
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(10.5)
    run2.font.color.rgb = TEXT_MUTED

# Workflow Diagram Box at bottom of left column
wf_card = add_card(slide5, Inches(1.0), Inches(5.4), Inches(5.1), Inches(1.1), bg_color=BADGE_BG, border_color=BADGE_TEXT)
tf_wf = wf_card.text_frame
tf_wf.margin_left = Inches(0.1)
tf_wf.margin_top = Inches(0.1)
p_wf = tf_wf.paragraphs[0]
p_wf.text = "CORE PATIENT WORKFLOW PIPELINE"
p_wf.font.size = Pt(10)
p_wf.font.bold = True
p_wf.font.color.rgb = BADGE_TEXT

p_wf2 = tf_wf.add_paragraph()
p_wf2.text = "Patient Registration -> Patient Record Saved -> Appointment Booked -> Hospital Care Workflow"
p_wf2.font.size = Pt(10)
p_wf2.font.bold = True
p_wf2.font.color.rgb = PRIMARY_NAVY

# Right Column: 3 Screenshot Placeholders
add_placeholder_frame(slide5, Inches(6.6), Inches(1.5), Inches(5.933), Inches(1.6), "INSERT SCREENSHOT - Patient Dashboard")
add_placeholder_frame(slide5, Inches(6.6), Inches(3.3), Inches(5.933), Inches(1.6), "INSERT SCREENSHOT - Patient Registration")
add_placeholder_frame(slide5, Inches(6.6), Inches(5.1), Inches(5.933), Inches(1.6), "INSERT SCREENSHOT - Patient List / Appointment Page")


# ==============================================================================
# SLIDE 6: MILESTONE 2 — CLINICAL & HOSPITAL MANAGEMENT MODULES
# ==============================================================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_header(slide6, "MILESTONE 2", "Milestone 2 - Clinical & Hospital Management Modules", "Comprehensive clinical modules for consultations, EHR, prescriptions, laboratory tests, billing, and notifications")
add_footer(slide6, 6)

# Top Left: Implemented Clinical Modules
add_card(slide6, Inches(0.8), Inches(1.5), Inches(6.0), Inches(4.1))
box6 = slide6.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.6), Inches(3.9))
tf6 = box6.text_frame
tf6.word_wrap = True

p = tf6.paragraphs[0]
p.text = "IMPLEMENTED CLINICAL MODULES"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

m2_items = [
    "Appointment Management: Scheduling, rescheduling, doctor assignment, and status updates ('Scheduled' -> 'Completed').",
    "Consultation Management: Recording symptoms, diagnosis, doctor notes, consultation fee, and auto-completing appointments.",
    "Electronic Health Records (EHR): Medical history, allergies, current medications, and vital signs (BP, Heart Rate, Temp, Weight).",
    "Prescription Management: Medication name, dosage, frequency, duration, special instructions, and prescriber tracking.",
    "Laboratory Management: Lab test requests, pending/completed status updates, results recording, and lab notes.",
    "Billing & Financial Module: Automated bill generation for consultations, lab tests, and pharmacy dispensing.",
    "Notification System: System alerts generated for low stock, new appointments, and bill updates."
]

for item in m2_items:
    p = tf6.add_paragraph()
    title, desc = item.split(": ", 1)
    run1 = p.add_run()
    run1.text = "• " + title + ": "
    run1.font.bold = True
    run1.font.size = Pt(9.5)
    run1.font.color.rgb = TEXT_DARK
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(9.5)
    run2.font.color.rgb = TEXT_MUTED

# Bottom Left: Module Relationship Flow Diagram
rel_card = add_card(slide6, Inches(0.8), Inches(5.7), Inches(6.0), Inches(1.0), bg_color=CARD_HDR_BG, border_color=PRIMARY_BLUE)
tf_rel = rel_card.text_frame
tf_rel.margin_left = Inches(0.15)
tf_rel.margin_top = Inches(0.1)

p = tf_rel.paragraphs[0]
p.text = "CLINICAL MODULE RELATIONSHIP FLOW"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = PRIMARY_BLUE

p = tf_rel.add_paragraph()
p.text = "Patient -> Appointment -> Consultation -> EHR / Prescription / Lab -> Billing / Notification"
p.font.size = Pt(10.5)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

# Right Column: 3 Screenshot Placeholders
add_placeholder_frame(slide6, Inches(7.1), Inches(1.5), Inches(5.433), Inches(1.6), "INSERT SCREENSHOT - Consultation & EHR Module")
add_placeholder_frame(slide6, Inches(7.1), Inches(3.3), Inches(5.433), Inches(1.6), "INSERT SCREENSHOT - Prescription & Lab Reports")
add_placeholder_frame(slide6, Inches(7.1), Inches(5.1), Inches(5.433), Inches(1.6), "INSERT SCREENSHOT - Billing & Patient Portal")


# ==============================================================================
# SLIDE 7: MILESTONE 3 — INTEGRATION, APIS & ADVANCED HOSPITAL OPERATIONS
# ==============================================================================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_background(slide7)
add_header(slide7, "MILESTONE 3", "Milestone 3 - Integration, APIs & Advanced Operations", "RESTful API suite, paginated endpoints, cross-module workflow integration, and global search")
add_footer(slide7, 7)

# Left Column: REST APIs & Advanced Features
add_card(slide7, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2))
box7 = slide7.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.4), Inches(5.0))
tf7 = box7.text_frame
tf7.word_wrap = True

p = tf7.paragraphs[0]
p.text = "REST APIS & ADVANCED INTEGRATIONS"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

m3_items = [
    "RESTful Endpoints: Implemented `/api/patients`, `/api/doctors`, `/api/consultations`, `/api/prescriptions`, `/api/lab_reports`, `/api/appointments`.",
    "API Authentication & RBAC: Session-protected API routes returning JSON responses for authorized clients.",
    "Query Pagination: Standardized pagination payload (`items`, `total`, `page`, `per_page`) for efficient frontend rendering.",
    "Appointment-Consultation Pipeline: Direct 'Start Consultation' quick action pre-populating patient & appointment context.",
    "Global Search Subsystem: Instant multi-entity search across patients, doctors, medical records, and bills.",
    "Integrated Notifications: Triggering real-time internal alerts upon appointment booking and bill updates."
]

for item in m3_items:
    p = tf7.add_paragraph()
    title, desc = item.split(": ", 1)
    run1 = p.add_run()
    run1.text = "• " + title + ": "
    run1.font.bold = True
    run1.font.size = Pt(10)
    run1.font.color.rgb = TEXT_DARK
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(10)
    run2.font.color.rgb = TEXT_MUTED

# REST API Workflow diagram at bottom of left box
api_wf_card = add_card(slide7, Inches(1.0), Inches(5.3), Inches(5.4), Inches(1.2), bg_color=BADGE_BG, border_color=BADGE_TEXT)
tf_api_wf = api_wf_card.text_frame
tf_api_wf.margin_left = Inches(0.1)
tf_api_wf.margin_top = Inches(0.08)

p = tf_api_wf.paragraphs[0]
p.text = "REST API DATA FLOW WORKFLOW"
p.font.size = Pt(9.5)
p.font.bold = True
p.font.color.rgb = BADGE_TEXT

p = tf_api_wf.add_paragraph()
p.text = "Client / Postman -> Flask REST API -> Auth/RBAC -> SQLAlchemy ORM -> MySQL -> JSON Response"
p.font.size = Pt(9.5)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

# Right Column: 2 Screenshot Placeholders
add_placeholder_frame(slide7, Inches(6.9), Inches(1.5), Inches(5.633), Inches(2.5), "INSERT SCREENSHOT - Postman API GET Request + JSON Response")
add_placeholder_frame(slide7, Inches(6.9), Inches(4.2), Inches(5.633), Inches(2.5), "INSERT SCREENSHOT - Integrated Workflow")


# ==============================================================================
# SLIDE 8: MILESTONE 4 — ANALYTICS, REPORTING, FEEDBACK, TESTING & DEPLOYMENT
# ==============================================================================
slide8 = prs.slides.add_slide(blank_layout)
set_slide_background(slide8)
add_header(slide8, "MILESTONE 4", "Milestone 4 - Analytics, Reporting, Feedback, Testing & Deployment", "Four core pillars: interactive analytics, administrative reporting, patient feedback, and production deployment")
add_footer(slide8, 8)

# 4 Quadrants / Cards Layout
cards_m4 = [
    ("A. ANALYTICS & DASHBOARDS", [
        "Real-time KPI metrics cards",
        "Monthly registration & appointment trends",
        "Disease distribution & lab stats charts",
        "Doctor-wise consultation counts & revenue",
        "In-memory caching for sub-15ms rendering"
    ], Inches(0.8), Inches(1.5)),
    
    ("B. ADMINISTRATIVE REPORTING", [
        "Specialized Patient & Consultation reports",
        "Doctor performance & monthly hospital reports",
        "PDF Report Generation via ReportLab engine",
        "Excel Workbook Export via XlsxWriter engine",
        "Streamed CSV exports & print support"
    ], Inches(6.8), Inches(1.5)),
    
    ("C. PATIENT FEEDBACK MODULE", [
        "1-5 star ratings (Doctor, Hospital, Lab, Pharmacy)",
        "Qualitative patient comments & notes",
        "Feedback history timeline & average calculation",
        "Admin feedback review & satisfaction metrics",
        "Exportable feedback summaries (CSV / PDF)"
    ], Inches(0.8), Inches(4.25)),
    
    ("D. TESTING & DEPLOYMENT", [
        "Functional, Security, API, UAT & Smoke testing",
        "30/30 Tests Passed (100% Pass Rate)",
        "Production WSGI server setup via Waitress",
        "Health check endpoint (`/health`)",
        "Database backup scripts (`backup_database.py`)"
    ], Inches(6.8), Inches(4.25))
]

for title, items, left, top in cards_m4:
    card = add_card(slide8, left, top, Inches(5.733), Inches(2.55))
    
    # Header strip inside card
    hdr = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.733), Inches(0.4))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = CARD_HDR_BG
    hdr.line.color.rgb = CARD_BORDER
    tf_h = hdr.text_frame
    tf_h.margin_left = Inches(0.15)
    tf_h.margin_top = Inches(0.08)
    p_h = tf_h.paragraphs[0]
    p_h.text = title
    p_h.font.size = Pt(11)
    p_h.font.bold = True
    p_h.font.color.rgb = PRIMARY_BLUE
    
    box = slide8.shapes.add_textbox(left + Inches(0.15), top + Inches(0.45), Inches(5.433), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK if "Passed" not in item else SUCCESS_GREEN
        if "Passed" in item:
            p.font.bold = True


# ==============================================================================
# SLIDE 9: COMPLETE SYSTEM WORKFLOW + MODULE OVERVIEW
# ==============================================================================
slide9 = prs.slides.add_slide(blank_layout)
set_slide_background(slide9)
add_header(slide9, "SYSTEM ARCHITECTURE", "End-to-End System Workflow & Module Overview", "Complete patient journey pipeline and overview of all 15 implemented hospital modules")
add_footer(slide9, 9)

# Top Box: End-to-End Workflow Diagram
add_card(slide9, Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.5))
wf_box = slide9.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.333), Inches(1.4))
tf_wf9 = wf_box.text_frame
tf_wf9.word_wrap = True

p = tf_wf9.paragraphs[0]
p.text = "END-TO-END CLINICAL WORKFLOW JOURNEY"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

flow_steps = [
    "1. Auth", "2. Patient Reg", "3. Appointment", "4. Consultation",
    "5. EHR Update", "6. Prescription", "7. Lab Test", "8. Billing",
    "9. Notification", "10. Feedback", "11. Analytics"
]
flow_str = " -> ".join(flow_steps)

p_flow = tf_wf9.add_paragraph()
p_flow.text = flow_str
p_flow.font.size = Pt(12)
p_flow.font.bold = True
p_flow.font.color.rgb = PRIMARY_BLUE

p_roles = tf_wf9.add_paragraph()
p_roles.text = "User Roles Involved: Admin (System/Reports), Doctor (Consultation/Rx), Nurse (Onboarding/Vitals), Pharmacist (Dispensing), Patient (Feedback/Portal)"
p_roles.font.size = Pt(10)
p_roles.font.color.rgb = TEXT_MUTED

# Bottom Box: 15 Implemented Modules Overview Grid
add_card(slide9, Inches(0.8), Inches(3.15), Inches(11.733), Inches(3.6))
mod_box = slide9.shapes.add_textbox(Inches(1.0), Inches(3.25), Inches(11.333), Inches(3.4))
tf_mod = mod_box.text_frame
tf_mod.word_wrap = True

p = tf_mod.paragraphs[0]
p.text = "OVERVIEW OF 15 IMPLEMENTED MODULES"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

modules = [
    ("1. Authentication & RBAC", "5 distinct portals, salted password hashing, audit logging"),
    ("2. Patient Management", "Demographic registration, Aadhaar validation, profile search"),
    ("3. Appointment Management", "Doctor scheduling, date/time pickers, status tracking"),
    ("4. Doctor Consultation", "Symptom recording, diagnosis, fee billing, appt completion"),
    ("5. Electronic Health Records", "Medical history, allergies, vital signs (BP, HR, Temp, Weight)"),
    ("6. Prescription Management", "Medication dosage, frequency, duration, prescriber attribution"),
    ("7. Laboratory Services", "Lab test requests, results recording, technician assignments"),
    ("8. Pharmacy & Dispensing", "Medicine inventory, batch tracking, dispensing billing"),
    ("9. Billing & Payments", "Automated invoice generation, payment status, method recording"),
    ("10. Notification System", "Internal system alerts for appointments, stock, and billing"),
    ("11. Global Search Engine", "Multi-table search across patients, doctors, records, and bills"),
    ("12. Administrative Reports", "Custom reports with PDF, Excel, and CSV multi-format exports"),
    ("13. Interactive Analytics", "Chart.js charts for registrations, diseases, revenue, and lab stats"),
    ("14. Patient Feedback", "1-5 star ratings for doctor/hospital/lab/pharmacy with comments"),
    ("15. RESTful API Suite", "Authenticated JSON endpoints with pagination support")
]

# Display modules in 3 columns
col_w = Inches(3.7)
for i, (m_title, m_desc) in enumerate(modules):
    col = i // 5
    row = i % 5
    m_left = Inches(1.0) + col * Inches(3.8)
    m_top = Inches(3.7) + row * Inches(0.58)
    
    m_card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, m_left, m_top, Inches(3.65), Inches(0.52))
    m_card.fill.solid()
    m_card.fill.fore_color.rgb = CARD_HDR_BG
    m_card.line.color.rgb = CARD_BORDER
    tf_mc = m_card.text_frame
    tf_mc.margin_left = Inches(0.08)
    tf_mc.margin_top = Inches(0.04)
    
    p = tf_mc.paragraphs[0]
    p.text = m_title
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    
    p_d = tf_mc.add_paragraph()
    p_d.text = m_desc
    p_d.font.size = Pt(8)
    p_d.font.color.rgb = TEXT_DARK


# ==============================================================================
# SLIDE 10: KEY FEATURES + APPLICATION SCREENSHOTS
# ==============================================================================
slide10 = prs.slides.add_slide(blank_layout)
set_slide_background(slide10)
add_header(slide10, "APPLICATION DEMO", "Key Features & Application Demonstration", "Core platform capability summary alongside UI application demonstration frames")
add_footer(slide10, 10)

# Left Column: Key Features List
add_card(slide10, Inches(0.8), Inches(1.5), Inches(5.2), Inches(5.2))
box10 = slide10.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(4.8), Inches(5.0))
tf10 = box10.text_frame
tf10.word_wrap = True

p = tf10.paragraphs[0]
p.text = "CORE SYSTEM CAPABILITIES"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

feats = [
    "Multi-Role Access Control (5 distinct user portals)",
    "Centralized Patient Demographic & Clinical Tracking",
    "Integrated Appointment-to-Consultation Pipeline",
    "Electronic Health Records (EHR) & Vitals Tracking",
    "Prescription, Lab Test & Pharmacy Integration",
    "Automated System Notifications & Alerts",
    "RESTful API Endpoints with Structured JSON Pagination",
    "Interactive Chart.js Analytics & Caching Engine",
    "Administrative Reporting with PDF/Excel/CSV Exports",
    "Patient Feedback & Satisfaction Rating System",
    "Security Hardening (Bcrypt, CSRF, SQLi, XSS Protection)",
    "Production WSGI Server & Health Endpoint Integration",
    "Automated MySQL Database Backup Scripting"
]

for feat in feats:
    p = tf10.add_paragraph()
    p.text = "+  " + feat
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK

# Right Column: 4 Large Screenshot Placeholders (2x2 Grid)
add_placeholder_frame(slide10, Inches(6.3), Inches(1.5), Inches(3.0), Inches(2.55), "SCREENSHOT 1 - Login/Dashboard")
add_placeholder_frame(slide10, Inches(9.5), Inches(1.5), Inches(3.0), Inches(2.55), "SCREENSHOT 2 - Patient/Appointment")
add_placeholder_frame(slide10, Inches(6.3), Inches(4.15), Inches(3.0), Inches(2.55), "SCREENSHOT 3 - Admin Analytics Dashboard")
add_placeholder_frame(slide10, Inches(9.5), Inches(4.15), Inches(3.0), Inches(2.55), "SCREENSHOT 4 - Reports/Feedback")


# ==============================================================================
# SLIDE 11: DATABASE DESIGN
# ==============================================================================
slide11 = prs.slides.add_slide(blank_layout)
set_slide_background(slide11)
add_header(slide11, "DATABASE ARCHITECTURE", "Database Design & Entity Relationships", "Relational database schema structure, SQLAlchemy ORM models, and foreign key relationships")
add_footer(slide11, 11)

# Left Box: Entity Relationships Text & Hierarchy
add_card(slide11, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.2))
box11 = slide11.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(6.1), Inches(5.0))
tf11 = box11.text_frame
tf11.word_wrap = True

p = tf11.paragraphs[0]
p.text = "IMPLEMENTED DATABASE ENTITIES & KEYS"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PRIMARY_NAVY

db_structure = [
    ("User Model (`users`)", "PK: `id` | Roles: Admin, Doctor, Nurse, Patient, Pharmacist"),
    ("Patient Model (`patients`)", "PK: `id` | Unique: `aadhaar` | Demographics & Disease"),
    ("Appointment Model (`appointments`)", "PK: `id` | FK: `patient_id` -> `patients.id`"),
    ("Consultation Model (`consultations`)", "PK: `id` | FK: `patient_id`, `doctor_id`, `bill_id`"),
    ("EHR Model (`ehrs`)", "PK: `id` | FK: `patient_id`, `doctor_id` | Vitals & Notes"),
    ("Prescription Model (`prescriptions`)", "PK: `id` | FK: `patient_id`, `doctor_id` | Medication & Dosage"),
    ("Lab Report Model (`lab_reports`)", "PK: `id` | FK: `patient_id`, `requested_by_id`, `performed_by_id`, `bill_id`"),
    ("Bill Model (`bills`)", "PK: `id` | FK: `patient_id` | Financial amounts & payment status"),
    ("Medicine & Dispensing (`medicines`, `dispensing_records`)", "FK: `patient_id`, `medicine_id`, `dispensed_by_id`, `bill_id`"),
    ("Notification & Feedback (`notifications`, `feedbacks`)", "FK: `user_id`, `patient_id`, `doctor_id`, `consultation_id`"),
    ("Login Activity (`login_activities`)", "FK: `user_id` | Audit status, action, timestamp, IP address")
]

for entity, keys in db_structure:
    p = tf11.add_paragraph()
    run1 = p.add_run()
    run1.text = "• " + entity + ": "
    run1.font.bold = True
    run1.font.size = Pt(9.5)
    run1.font.color.rgb = TEXT_DARK
    run2 = p.add_run()
    run2.text = keys
    run2.font.size = Pt(9)
    run2.font.color.rgb = TEXT_MUTED

# Right Column: ER Diagram Placeholder
add_placeholder_frame(slide11, Inches(7.5), Inches(1.5), Inches(5.033), Inches(5.2), "INSERT DATABASE / ER DIAGRAM SCREENSHOT IF REQUIRED")


# ==============================================================================
# SLIDE 12: TESTING, SECURITY, PERFORMANCE & DEPLOYMENT
# ==============================================================================
slide12 = prs.slides.add_slide(blank_layout)
set_slide_background(slide12)
add_header(slide12, "QUALITY & DEPLOYMENT", "Testing, Security, Performance & Production Readiness", "Verification of automated unit test suite, security hardening, performance benchmarks, and server config")
add_footer(slide12, 12)

# 4 Quadrants Layout for Testing, Security, Performance, Deployment
cards_m12 = [
    ("AUTOMATED TESTING SUITE", [
        "6 Test Suites: Smoke, Functional, Security, Performance, UAT, Feedback",
        "Verified Results: 30 Tests / 30 Passed / 0 Failed",
        "100% Test Pass Rate across complete application",
        "Automated UAT verifying end-to-end patient care journey"
    ], Inches(0.8), Inches(1.5), SUCCESS_GREEN),
    
    ("SECURITY HARDENING", [
        "Werkzeug / Bcrypt salted password hashing",
        "Role-Based Access Control (RBAC) portal enforcement",
        "CSRF protection on form submissions via Flask-WTF",
        "SQL Injection prevention via SQLAlchemy parameterized queries",
        "Jinja2 auto-escaping protecting against XSS attacks"
    ], Inches(6.8), Inches(1.5), PRIMARY_NAVY),
    
    ("PERFORMANCE OPTIMIZATION", [
        "In-Memory Dashboard Caching: Render time ~12 ms (vs ~32 ms uncached)",
        "REST API Latency: ~2.7 ms (/api/patients), ~1.7 ms (/api/doctors)",
        "Database Indexing: Indexed role, status, and timestamp fields",
        "Structured API pagination reducing payload size"
    ], Inches(0.8), Inches(4.15), PRIMARY_BLUE),
    
    ("PRODUCTION DEPLOYMENT & HEALTH", [
        "Environment variables management via python-dotenv (.env)",
        "Production WSGI Server: Multi-threaded Waitress WSGI server",
        "Health Check Endpoint: GET /health returning JSON status",
        "Database Backup: Timestamped SQL dump via backup_database.py"
    ], Inches(6.8), Inches(4.15), TEAL_ACCENT)
]

for title, items, left, top, accent_color in cards_m12:
    card = add_card(slide12, left, top, Inches(5.733), Inches(2.65))
    
    hdr = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.733), Inches(0.4))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = CARD_HDR_BG
    hdr.line.color.rgb = CARD_BORDER
    tf_h = hdr.text_frame
    tf_h.margin_left = Inches(0.15)
    tf_h.margin_top = Inches(0.08)
    p_h = tf_h.paragraphs[0]
    p_h.text = title
    p_h.font.size = Pt(11)
    p_h.font.bold = True
    p_h.font.color.rgb = accent_color
    
    box = slide12.shapes.add_textbox(left + Inches(0.15), top + Inches(0.45), Inches(5.433), Inches(2.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_DARK
        if "30 Tests / 30 Passed" in item:
            p.font.bold = True
            p.font.color.rgb = SUCCESS_GREEN

# Small screenshot placeholders overlaid/positioned cleanly
add_placeholder_frame(slide12, Inches(4.2), Inches(3.2), Inches(2.2), Inches(0.85), "TERMINAL - 30/30 PASSED")
add_placeholder_frame(slide12, Inches(10.2), Inches(5.8), Inches(2.2), Inches(0.85), "/health RESPONSE")


# ==============================================================================
# SLIDE 13: CHALLENGES, LEARNING OUTCOMES, FUTURE SCOPE & CONCLUSION
# ==============================================================================
slide13 = prs.slides.add_slide(blank_layout)
set_slide_background(slide13)
add_header(slide13, "FUTURE SCOPE & SUMMARY", "Challenges, Learning Outcomes, Future Scope & Conclusion", "Summary of implementation learnings, potential enhancements, and final evaluation conclusions")
add_footer(slide13, 13)

cards_m13 = [
    ("CHALLENGES OVERCOME", [
        "Cross-Module Relational Integrity: Maintaining FK constraints across consultations, prescriptions, labs, and bills.",
        "Role Authorization Enforcement: Implementing strict RBAC checks across 5 distinct portals.",
        "Reporting Engine Design: Generating clean PDF/Excel layouts from complex SQL queries."
    ], Inches(0.8), Inches(1.5)),
    
    ("KEY LEARNING OUTCOMES", [
        "Full-Stack Flask Architecture: Building modular web applications with SQLAlchemy ORM.",
        "RESTful API & Serialization: Designing paginated JSON endpoints and authentication.",
        "Quality Assurance: Writing comprehensive automated unit and UAT test suites."
    ], Inches(6.8), Inches(1.5)),
    
    ("FUTURE ENHANCEMENTS (PLANNED)", [
        "Cloud Platform Deployment: Hosting on AWS / Azure with automated cloud backups.",
        "Email & SMS Gateway Integration: Automated appointment reminders and alert notifications.",
        "Predictive AI & Analytics: Machine learning models for clinical risk prediction."
    ], Inches(0.8), Inches(4.15)),
    
    ("CONCLUSION & PROJECT SUMMARY", [
        "IPCMS integrates patient care, hospital operations, analytics, reporting, and deployment into one robust platform.",
        "Validated with 100% automated test pass rate and ready for Waitress WSGI server deployment.",
        "Thank You  |  Questions & Discussion  |  Contact: [Insert Email]"
    ], Inches(6.8), Inches(4.15))
]

for title, items, left, top in cards_m13:
    card = add_card(slide13, left, top, Inches(5.733), Inches(2.65))
    
    hdr = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.733), Inches(0.4))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = CARD_HDR_BG
    hdr.line.color.rgb = CARD_BORDER
    tf_h = hdr.text_frame
    tf_h.margin_left = Inches(0.15)
    tf_h.margin_top = Inches(0.08)
    p_h = tf_h.paragraphs[0]
    p_h.text = title
    p_h.font.size = Pt(11)
    p_h.font.bold = True
    p_h.font.color.rgb = PRIMARY_BLUE
    
    box = slide13.shapes.add_textbox(left + Inches(0.15), top + Inches(0.45), Inches(5.433), Inches(2.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if ": " in item:
            title, desc = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = "• " + title + ": "
            run1.font.bold = True
            run1.font.size = Pt(9.5)
            run1.font.color.rgb = PRIMARY_BLUE
            run2 = p.add_run()
            run2.text = desc
            run2.font.size = Pt(9.5)
            run2.font.color.rgb = TEXT_DARK
        else:
            p.text = "• " + item
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_DARK
            if "Thank You" in item:
                p.font.bold = True
                p.font.color.rgb = PRIMARY_NAVY

# Save presentation
output_path = "IPCMS_Project_Evaluation_Milestone_1_to_4.pptx"
prs.save(output_path)
print(f"Presentation saved successfully as '{output_path}'. Total slides: {len(prs.slides)}")
