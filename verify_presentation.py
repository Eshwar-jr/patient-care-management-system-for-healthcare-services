import sys
from pptx import Presentation

prs = Presentation("IPCMS_Project_Evaluation_Milestone_1_to_4.pptx")
print(f"Total Slides: {len(prs.slides)}")

assert len(prs.slides) == 13, f"Expected 13 slides, got {len(prs.slides)}"

titles = []
for i, slide in enumerate(prs.slides):
    slide_title = ""
    # Find text boxes or shapes with title
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if "Milestone" in text or "Integrated" in text or "Problem" in text or "Technology" in text or "Workflow" in text or "Key Features" in text or "Database" in text or "Testing" in text or "Challenges" in text or "End-to-End" in text:
                    if not slide_title:
                        slide_title = text
    titles.append((i + 1, slide_title))
    print(f"Slide {i+1}: Shapes count = {len(slide.shapes)}")

print("\n--- Slide Titles Breakdown ---")
for num, title in titles:
    print(f"Slide {num}: {title}")

# Verify milestone slide counts
m1_slides = [t for n, t in titles if "Milestone 1" in t]
m2_slides = [t for n, t in titles if "Milestone 2" in t]
m3_slides = [t for n, t in titles if "Milestone 3" in t]
m4_slides = [t for n, t in titles if "Milestone 4" in t]

print(f"\nMilestone 1 Slides Count: {len(m1_slides)} (Expected 2)")
print(f"Milestone 2 Slides Count: {len(m2_slides)} (Expected 1)")
print(f"Milestone 3 Slides Count: {len(m3_slides)} (Expected 1)")
print(f"Milestone 4 Slides Count: {len(m4_slides)} (Expected 1)")

assert len(m1_slides) == 2, "Milestone 1 must have exactly 2 slides"
assert len(m2_slides) == 1, "Milestone 2 must have exactly 1 slide"
assert len(m3_slides) == 1, "Milestone 3 must have exactly 1 slide"
assert len(m4_slides) == 1, "Milestone 4 must have exactly 1 slide"

print("\n[OK] Presentation verification successful!")
